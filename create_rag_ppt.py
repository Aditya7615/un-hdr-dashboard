from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── Helper functions ───────────────────────────────────────────────
def add_title_slide(prs, title, subtitle, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(2))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0x4A, 0x4A, 0x6A)
        p2.alignment = PP_ALIGN.CENTER

    if note:
        txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.8))
        tf2 = txb2.text_frame
        p3 = tf2.paragraphs[0]
        p3.text = note
        p3.font.size = Pt(14)
        p3.font.italic = True
        p3.font.color.rgb = RGBColor(0x88, 0x88, 0x99)
        p3.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.33), Inches(1.5))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.9))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Bullets
    txb2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.5))
    tf2 = txb2.text_frame
    tf2.word_wrap = True

    for i, (bold_part, rest) in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"• {bold_part}{rest}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x2D, 0x2D, 0x44)
        p.space_after = Pt(10)

    if note:
        txb3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.6))
        tf3 = txb3.text_frame
        pn = tf3.paragraphs[0]
        pn.text = f"📌 {note}"
        pn.font.size = Pt(13)
        pn.font.italic = True
        pn.font.color.rgb = RGBColor(0x88, 0x88, 0x99)
    return slide

def add_two_column_slide(prs, title, left_header, left_bullets, right_header, right_bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.9))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Left column header
    txb_l = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.5))
    tf_l = txb_l.text_frame
    ph = tf_l.paragraphs[0]
    ph.text = left_header
    ph.font.size = Pt(22)
    ph.font.bold = True
    ph.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Left bullets
    txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(5.8), Inches(4.5))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    for i, b in enumerate(left_bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x2D, 0x2D, 0x44)
        p.space_after = Pt(8)

    # Right column header
    txb_r = slide.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(5.8), Inches(0.5))
    tf_r = txb_r.text_frame
    phr = tf_r.paragraphs[0]
    phr.text = right_header
    phr.font.size = Pt(22)
    phr.font.bold = True
    phr.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Right bullets
    txb3 = slide.shapes.add_textbox(Inches(6.8), Inches(1.75), Inches(5.8), Inches(4.5))
    tf3 = txb3.text_frame
    tf3.word_wrap = True
    for i, b in enumerate(right_bullets):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x2D, 0x2D, 0x44)
        p.space_after = Pt(8)

    if note:
        txb4 = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.6))
        tf4 = txb4.text_frame
        pn = tf4.paragraphs[0]
        pn.text = f"📌 {note}"
        pn.font.size = Pt(13)
        pn.font.italic = True
        pn.font.color.rgb = RGBColor(0x88, 0x88, 0x99)
    return slide

# ─── Slide 1: Title ──────────────────────────────────────────────────
add_title_slide(
    prs,
    "Retrieval-Augmented Generation (RAG)",
    "Enhancing LLMs with External Knowledge & Performance Analysis",
    "From 'Generative AI' → 'Grounded AI'"
)

# ─── Slide 2: Why RAG ────────────────────────────────────────────────
add_content_slide(
    prs,
    "The 'Why' Behind RAG",
    [
        ("Knowledge Cut-offs: ", "Training data is static — LLMs don't know about events after their training date."),
        ("Hallucinations: ", "Models confidently generate factually incorrect statements."),
        ("Data Privacy: ", "Standard models have no access to your private/internal documents."),
        ("The RAG Solution: ", "Connects the LLM to a live 'library' of your own data — no fine-tuning required."),
    ],
    note="RAG = Retrieval + Generation. The model stays grounded in real data."
)

# ─── Slide 3: Architecture Overview ─────────────────────────────────
add_content_slide(
    prs,
    "High-Level RAG Architecture",
    [
        ("Ingestion: ", "Preparing & loading data into the knowledge base (chunking, embedding, indexing)."),
        ("Retrieval: ", "Converting the user query to a vector and searching for the most relevant chunks."),
        ("Generation: ", "Feeding retrieved context + question to the LLM for grounded answer generation."),
    ],
    note="Query → Vector DB (similarity search) → Top-k Chunks → LLM Prompt → Answer"
)

# ─── Slide 4: Phase 1 – Data Ingestion ─────────────────────────────
add_content_slide(
    prs,
    "Phase 1: Data Ingestion (The Preparation)",
    [
        ("Chunking: ", "Splitting long documents into smaller, semantically meaningful pieces."),
        ("Embeddings: ", "Using an AI model (e.g., OpenAI ada, BGE) to convert text into N-dimensional vectors."),
        ("Vector Database: ", "Specialized storage (Pinecone, Weaviate, Chroma, FAISS) for fast similarity search."),
        ("Indexing: ", "Building the search index over embedded chunks for sub-second retrieval."),
    ],
    note="Semantic Search finds meaning, not just keyword matches."
)

# ─── Slide 5: Phase 2 – RAG Workflow ───────────────────────────────
add_content_slide(
    prs,
    "Phase 2: The RAG Workflow (The Execution)",
    [
        ("Query Transformation: ", "The user's question is converted into a vector using the same embedding model."),
        ("Similarity Search: ", "Top-k most relevant chunks are pulled from the vector database."),
        ("Prompt Augmentation: ", "Retrieved chunks are injected into the LLM prompt as 'Context'."),
        ("Generation: ", "The LLM produces an answer strictly based on the provided context."),
    ],
    note="Zero hallucinations (in theory) — the model can only 'see' the context you give it."
)

# ─── Slide 6: RAGAS Framework ──────────────────────────────────────
add_content_slide(
    prs,
    "Performance Analysis: The RAGAS Framework",
    [
        ("Context Precision: ", "Did we retrieve the right documents among all relevant ones?"),
        ("Context Recall: ", "Did we retrieve enough of the relevant information to fully answer?"),
        ("Faithfulness: ", "Did the LLM stick strictly to the provided context, or did it hallucinate?"),
        ("Answer Relevance: ", "Is the final response actually useful and directly responsive to the query?"),
    ],
    note="RAGAS = RAG Assessment. A framework for measuring RAG pipeline quality quantitatively."
)

# ─── Slide 7: Technical Metrics ────────────────────────────────────
add_two_column_slide(
    prs,
    "Technical Performance Metrics",
    "Speed & Cost",
    [
        "Latency: RAG adds retrieval time — Time-to-First-Token (TTFT) is higher than standard chat.",
        "Token Efficiency: Balancing context window size vs. API call cost per query.",
        "Indexing Cost: One-time upfront cost of embedding and storing the knowledge base.",
    ],
    "Quality & Accuracy",
    [
        "Hit Rate: How often the correct answer exists in our top-k retrieved results.",
        "MRR (Mean Reciprocal Rank): Is the best answer ranked #1 on average?",
        "NDCG: Normalized Discounted Cumulative Gain — how well results are ordered.",
    ],
    note="Trade-off: More context → better answers but higher latency and cost."
)

# ─── Slide 8: Challenges & Best Practices ──────────────────────────
add_two_column_slide(
    prs,
    "Challenges & Best Practices",
    "Common Pitfalls",
    [
        "Poor Chunking: Cutting mid-sentence loses semantic meaning.",
        "Retrieval Noise: Irrelevant chunks pollute context and confuse the LLM.",
        "Embedding Drift: Chunks lose relevance as data updates over time.",
        "Cold-start: New/empty vector DBs return poor results initially.",
    ],
    "Best Practices",
    [
        "Re-ranking models (e.g., Cohere Rerank) refine search results after initial retrieval.",
        "Hybrid Search: Combine keyword (BM25) + semantic (vector) search.",
        "Metadata Filtering: Narrow retrieval by date, source, or category.",
        "Chunk Overlap: 10–20% overlap preserves context across chunk boundaries.",
    ]
)

# ─── Slide 9: Conclusion ────────────────────────────────────────────
add_content_slide(
    prs,
    "Conclusion",
    [
        ("RAG = Reasoning + Factual Grounding: ", "Bridges the gap between a model's reasoning ability and real-time data."),
        ("Success Formula: ", "High Retriever Accuracy × High Generator Faithfulness = Reliable RAG system."),
        ("Fail-safe: ", "A well-built RAG system should say 'I don't know' rather than hallucinate."),
        ("Key Takeaway: ", "RAG makes LLMs auditable, up-to-date, and domain-specific — without retraining."),
    ],
    note="The best RAG system is one that knows what it doesn't know."
)

# ─── Slide 10: Q&A ──────────────────────────────────────────────────
add_title_slide(
    prs,
    "Q&A",
    "Questions on RAG Architecture, Workflow, or Performance Metrics?",
    ""
)

# ─── Save ──────────────────────────────────────────────────────────
output_path = "/Users/adityagoyal/work/Lengraph/RAG_Presentation.pptx"
prs.save(output_path)
print(f"Saved → {output_path}")
